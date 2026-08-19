"""Extracts usable content from an uploaded file so the pipeline can build a workflow
from real customer-provided context instead of typed text or synthesized samples.

Two outcomes:
  - "tabular" (.csv, .xlsx/.xls): the file's own rows become the REAL dataset the
    Executor runs against - no synthetic data is generated at all for these.
  - "text" (.docx, .pdf, .pptx): extracted text is treated exactly like a typed
    problem description and fed into the Requirement Agent; the sample dataset is
    still synthesized afterward, same as typing a prompt (there's no tabular ground
    truth to run against).
"""
from __future__ import annotations

import csv
import io
import re
from dataclasses import dataclass, field

MAX_TABULAR_ROWS = 200
MAX_TEXT_CHARS = 8000
MAX_PDF_PAGES = 30

_ALLOWED_EXTENSIONS = {"csv", "xlsx", "xls", "docx", "pdf", "pptx"}


class UnsupportedFileTypeError(Exception):
    pass


class FileParseError(Exception):
    pass


@dataclass
class ExtractedContent:
    kind: str  # "tabular" | "text"
    text: str = ""
    records: list[dict] = field(default_factory=list)
    fields: list[str] = field(default_factory=list)
    record_label: str = "record"


def _record_label_from_filename(filename: str) -> str:
    stem = re.sub(r"\.[^.]+$", "", filename)
    stem = re.sub(r"[_\-]+", " ", stem).strip().lower() or "record"
    if stem.endswith("ies"):
        stem = stem[:-3] + "y"
    elif stem.endswith("ses"):
        stem = stem[:-2]
    elif stem.endswith("s") and not stem.endswith("ss"):
        stem = stem[:-1]
    return stem or "record"


def _coerce(value: str | None) -> object:
    """csv.DictReader hands back every cell as a string - without this, a numeric
    column like risk_score="72" silently fails every CHECK_CONDITION/ANALYZE
    comparison (TypeError comparing str to int, caught and treated as unmatched)."""
    if value is None:
        return None
    stripped = value.strip()
    if stripped == "":
        return ""
    if re.fullmatch(r"-?\d+", stripped):
        return int(stripped)
    if re.fullmatch(r"-?\d*\.\d+", stripped):
        return float(stripped)
    if stripped.lower() in ("true", "false"):
        return stripped.lower() == "true"
    return value


def _extract_csv(content: bytes, filename: str) -> ExtractedContent:
    try:
        text = content.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = content.decode("latin-1")
    reader = csv.DictReader(io.StringIO(text))
    fields = list(reader.fieldnames or [])
    records = []
    for i, row in enumerate(reader):
        if i >= MAX_TABULAR_ROWS:
            break
        records.append({k: _coerce(v) for k, v in row.items()})
    if not fields or not records:
        raise FileParseError("The CSV file has no header row or no data rows.")
    return ExtractedContent(kind="tabular", records=records, fields=fields, record_label=_record_label_from_filename(filename))


def _extract_xlsx(content: bytes, filename: str) -> ExtractedContent:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sheet = wb.worksheets[0]
    rows_iter = sheet.iter_rows(values_only=True)
    try:
        header = [str(h) if h is not None else f"column_{i}" for i, h in enumerate(next(rows_iter))]
    except StopIteration:
        raise FileParseError("The spreadsheet's first sheet is empty.")

    records = []
    for i, row in enumerate(rows_iter):
        if i >= MAX_TABULAR_ROWS:
            break
        records.append({header[j]: row[j] for j in range(min(len(header), len(row)))})
    if not records:
        raise FileParseError("The spreadsheet's first sheet has a header row but no data rows.")
    return ExtractedContent(kind="tabular", records=records, fields=header, record_label=_record_label_from_filename(filename))


def _extract_docx(content: bytes) -> ExtractedContent:
    from docx import Document

    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            paragraphs.append(" | ".join(cell.text for cell in row.cells))
    text = "\n".join(paragraphs)[:MAX_TEXT_CHARS]
    if not text.strip():
        raise FileParseError("No readable text found in the Word document.")
    return ExtractedContent(kind="text", text=text)


def _extract_pdf(content: bytes) -> ExtractedContent:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages_text = []
    for page in reader.pages[:MAX_PDF_PAGES]:
        pages_text.append(page.extract_text() or "")
    text = "\n".join(pages_text)[:MAX_TEXT_CHARS]
    if not text.strip():
        raise FileParseError("No extractable text found in the PDF (it may be a scanned/image-only PDF).")
    return ExtractedContent(kind="text", text=text)


def _extract_pptx(content: bytes) -> ExtractedContent:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text)
    text = "\n".join(chunks)[:MAX_TEXT_CHARS]
    if not text.strip():
        raise FileParseError("No readable text found in the PowerPoint file.")
    return ExtractedContent(kind="text", text=text)


def extract_from_file(filename: str, content: bytes) -> ExtractedContent:
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in _ALLOWED_EXTENSIONS:
        raise UnsupportedFileTypeError(
            f"'.{ext}' is not supported. Allowed: {', '.join(sorted(_ALLOWED_EXTENSIONS))}."
        )
    if ext == "csv":
        return _extract_csv(content, filename)
    if ext in ("xlsx", "xls"):
        return _extract_xlsx(content, filename)
    if ext == "docx":
        return _extract_docx(content)
    if ext == "pdf":
        return _extract_pdf(content)
    if ext == "pptx":
        return _extract_pptx(content)
    raise UnsupportedFileTypeError(ext)  # unreachable, kept for exhaustiveness
